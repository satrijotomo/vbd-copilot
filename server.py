"""FastAPI server exposing CSA-Copilot over HTTP + WebSocket.

Start with:
    python app.py --server [--port PORT]

The server picks a random free port if ``--port`` is omitted and writes the
chosen port to stdout as ``PORT:XXXX`` so the Electron main process can read it.

Security:
  - Binds only to 127.0.0.1 (never 0.0.0.0).
  - The /file endpoint validates all paths are under ``outputs/``.
  - WebSocket messages are validated before processing.
  - No secrets are exposed through any endpoint.
"""

from __future__ import annotations

import asyncio
import base64
import contextlib
import io
import json
import logging
import os
import time
from pathlib import Path
from typing import Any

from fastapi import FastAPI, HTTPException, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# App-wide state injected at startup by server_main()
# ---------------------------------------------------------------------------

_app_dir: Path = Path(__file__).resolve().parent
_outputs_dir: Path = _app_dir / "outputs"
_event_store: Any = None        # EventStore instance
_copilot_client: Any = None     # CopilotClient instance
_session_map: dict[str, Any] = {}   # session_id -> Session object
_collector: Any = None          # EventCollector instance

# ---------------------------------------------------------------------------
# FastAPI application
# ---------------------------------------------------------------------------

app = FastAPI(title="CSA Copilot API", version="1.0.0", docs_url=None, redoc_url=None)

# Allow the Electron renderer (loaded from file://) to call the API.
# Restricted to only localhost origins for security.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost", "http://127.0.0.1"],
    allow_credentials=False,
    allow_methods=["GET", "POST", "DELETE"],
    allow_headers=["Content-Type"],
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _store() -> Any:
    if _event_store is None:
        raise HTTPException(status_code=503, detail="Store not initialised")
    return _event_store


def _safe_outputs_path(raw: str) -> Path:
    """Resolve *raw* and assert it is inside outputs/.

    Raises HTTPException 400 if the path escapes the outputs directory.
    """
    try:
        resolved = Path(raw).resolve()
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid path")
    outputs_resolved = _outputs_dir.resolve()
    try:
        resolved.relative_to(outputs_resolved)
    except ValueError:
        raise HTTPException(status_code=400, detail="Path outside outputs directory")
    return resolved


# ---------------------------------------------------------------------------
# Health
# ---------------------------------------------------------------------------

@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok", "version": "1.0.0"}


# ---------------------------------------------------------------------------
# Agents
# ---------------------------------------------------------------------------

@app.get("/agents")
async def list_agents() -> JSONResponse:
    from agents import CATALOG
    result = []
    for name, agent in CATALOG.all_agents.items():
        result.append({
            "name": name,
            "display_name": agent.display_name,
            "description": agent.description,
            "model": getattr(agent, "model", ""),
            "infer": agent.infer,
        })
    return JSONResponse(content=result)


# ---------------------------------------------------------------------------
# Sessions
# ---------------------------------------------------------------------------

class CreateSessionRequest(BaseModel):
    agent: str | None = None
    model: str | None = None


@app.post("/sessions")
async def create_session(body: CreateSessionRequest) -> JSONResponse:
    from agents import ALL_AGENT_CONFIGS, ALL_SKILL_DIRS, DEFAULT_MODEL
    from tools import ALL_CUSTOM_TOOLS

    if _copilot_client is None:
        raise HTTPException(status_code=503, detail="Copilot client not ready")

    model = body.model or DEFAULT_MODEL

    async def _perm(request: Any, _inv: Any) -> Any:
        from copilot.types import PermissionRequestResult
        return PermissionRequestResult(kind="approved")

    async def _user_input(request: Any, _inv: Any) -> Any:
        from copilot.types import UserInputResponse
        from server_adapter import pop_user_response
        question = request.get("question", "")
        choices = request.get("choices")
        from server_adapter import set_active_ws, _send
        _send({"type": "waiting_for_input", "question": question, "choices": choices})
        try:
            answer = await pop_user_response(timeout=300.0)
        except asyncio.TimeoutError:
            answer = ""
        return UserInputResponse(answer=answer, wasFreeform=True)

    try:
        session = await _copilot_client.create_session(
            {
                "model": model,
                "streaming": True,
                "custom_agents": ALL_AGENT_CONFIGS,
                "tools": ALL_CUSTOM_TOOLS,
                "skill_directories": ALL_SKILL_DIRS,
                "on_permission_request": _perm,
                "on_user_input_request": _user_input,
                "working_directory": str(_app_dir),
            }
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))

    sid = session.session_id
    _session_map[sid] = session

    if _collector:
        _collector.on_session_created(
            sid, agent=body.agent or "", model=model, frontend="desktop"
        )
    if _event_store and body.agent:
        _event_store.update_session_agent(sid, body.agent)

    return JSONResponse(content={"session_id": sid, "model": model})


class ResumeSessionRequest(BaseModel):
    pass


@app.post("/sessions/{session_id}/resume")
async def resume_session(session_id: str, body: ResumeSessionRequest) -> JSONResponse:
    from agents import ALL_AGENT_CONFIGS, ALL_SKILL_DIRS, DEFAULT_MODEL
    from tools import ALL_CUSTOM_TOOLS

    if _copilot_client is None:
        raise HTTPException(status_code=503, detail="Copilot client not ready")
    if _event_store is None:
        raise HTTPException(status_code=503, detail="Store not initialised")

    full_id = _event_store.resolve_prefix("sessions", session_id)
    if not full_id:
        raise HTTPException(status_code=404, detail="Session not found")

    s_detail = _event_store.get_session(full_id)
    if not s_detail or not s_detail.get("resumable"):
        raise HTTPException(status_code=400, detail="Session is not resumable")

    async def _perm(request: Any, _inv: Any) -> Any:
        from copilot.types import PermissionRequestResult
        return PermissionRequestResult(kind="approved")

    async def _user_input(request: Any, _inv: Any) -> Any:
        from copilot.types import UserInputResponse
        from server_adapter import _send, pop_user_response
        question = request.get("question", "")
        choices = request.get("choices")
        _send({"type": "waiting_for_input", "question": question, "choices": choices})
        try:
            answer = await pop_user_response(timeout=300.0)
        except asyncio.TimeoutError:
            answer = ""
        return UserInputResponse(answer=answer, wasFreeform=True)

    try:
        session = await _copilot_client.resume_session(
            full_id,
            {
                "streaming": True,
                "custom_agents": ALL_AGENT_CONFIGS,
                "tools": ALL_CUSTOM_TOOLS,
                "skill_directories": ALL_SKILL_DIRS,
                "on_permission_request": _perm,
                "on_user_input_request": _user_input,
                "working_directory": str(_app_dir),
            },
        )
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))

    _session_map[full_id] = session
    _event_store.reactivate_session(full_id)

    s_detail = _event_store.get_session(full_id)
    turns = _event_store.get_turns(full_id)

    return JSONResponse(content={
        "session_id": full_id,
        "agent": s_detail.get("agent", ""),
        "model": s_detail.get("model", ""),
        "turn_count": s_detail.get("turn_count", 0),
        "turns": turns,
    })


@app.delete("/sessions/{session_id}")
async def end_session(session_id: str) -> JSONResponse:
    if _event_store:
        full_id = _event_store.resolve_prefix("sessions", session_id) or session_id
        _event_store.end_session(full_id, resumable=False)
    session = _session_map.pop(session_id, None)
    if session:
        with contextlib.suppress(Exception):
            session._event_handlers.clear()
    return JSONResponse(content={"ok": True})


@app.get("/sessions")
async def list_sessions(all: bool = False) -> JSONResponse:
    store = _store()
    sessions = store.list_sessions(include_ended=all)
    return JSONResponse(content=sessions)


@app.get("/sessions/{session_id}")
async def get_session(session_id: str) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    detail = store.get_session(full_id)
    if not detail:
        raise HTTPException(status_code=404, detail="Session not found")
    return JSONResponse(content=detail)


@app.get("/sessions/{session_id}/turns")
async def get_session_turns(session_id: str) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    turns = store.get_turns(full_id)
    return JSONResponse(content=turns)


@app.get("/sessions/{session_id}/turns/{turn_number}/invocations")
async def get_turn_invocations(session_id: str, turn_number: int) -> JSONResponse:
    store = _store()
    full_id = store.resolve_prefix("sessions", session_id) or session_id
    turns = store.get_turns(full_id)
    turn = next((t for t in turns if t.get("turn_number") == turn_number), None)
    if not turn:
        raise HTTPException(status_code=404, detail="Turn not found")
    invocations = store.get_invocations(turn["id"])
    return JSONResponse(content=invocations)


# ---------------------------------------------------------------------------
# Usage
# ---------------------------------------------------------------------------

@app.get("/usage")
async def get_usage(
    period: str = "all",
    agent: str | None = None,
    model: str | None = None,
) -> JSONResponse:
    store = _store()
    data = store.get_usage_stats(period=period, agent=agent, model=model)
    return JSONResponse(content=data)


# ---------------------------------------------------------------------------
# Outputs
# ---------------------------------------------------------------------------

_FILE_TYPE_MAP = {
    ".pptx": "pptx",
    ".md": "markdown",
    ".py": "python",
    ".sh": "shell",
    ".bicep": "bicep",
    ".json": "json",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".ts": "typescript",
    ".txt": "text",
}

_SKIP_DIRS = {".fragments", "__pycache__", ".git", "node_modules"}


def _classify_output_category(path: Path) -> str:
    parts = path.parts
    for i, part in enumerate(parts):
        if part == "outputs" and i + 1 < len(parts):
            return parts[i + 1]
    return "other"


@app.get("/outputs")
async def list_outputs() -> JSONResponse:
    outputs_resolved = _outputs_dir.resolve()
    results = []
    for p in outputs_resolved.rglob("*"):
        if not p.is_file():
            continue
        if any(part in _SKIP_DIRS for part in p.parts):
            continue
        suffix = p.suffix.lower()
        if suffix not in _FILE_TYPE_MAP:
            continue
        if "-plan.md" in p.name:
            continue
        try:
            stat = p.stat()
        except OSError:
            continue
        rel = str(p.relative_to(outputs_resolved))
        results.append({
            "path": str(p),
            "relative": rel,
            "name": p.name,
            "type": _FILE_TYPE_MAP.get(suffix, "file"),
            "category": _classify_output_category(p),
            "size": stat.st_size,
            "modified": stat.st_mtime,
        })
    results.sort(key=lambda x: x["modified"], reverse=True)
    return JSONResponse(content=results)


@app.get("/file")
async def read_file(path: str) -> JSONResponse:
    resolved = _safe_outputs_path(path)
    if not resolved.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    try:
        content = resolved.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc))
    return JSONResponse(content={"path": str(resolved), "content": content})


# ---------------------------------------------------------------------------
# PPTX Preview
# ---------------------------------------------------------------------------

class PptxPreviewRequest(BaseModel):
    path: str
    max_width: int = 1280


@app.post("/preview/pptx")
async def preview_pptx(body: PptxPreviewRequest) -> JSONResponse:
    resolved = _safe_outputs_path(body.path)
    if not resolved.is_file() or resolved.suffix.lower() != ".pptx":
        raise HTTPException(status_code=400, detail="Not a .pptx file")

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed")

    try:
        from PIL import Image
    except ImportError:
        raise HTTPException(status_code=500, detail="Pillow not installed")

    try:
        prs = Presentation(str(resolved))
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Cannot open PPTX: {exc}")

    slides_data = []
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    # Calculate pixel dimensions preserving aspect ratio
    px_width = body.max_width
    aspect = slide_height_emu / slide_width_emu if slide_width_emu else 9 / 16
    px_height = int(px_width * aspect)

    for i, slide in enumerate(prs.slides):
        # Build a white canvas with text overlay using Pillow
        # (Full rendering requires LibreOffice; we create a styled placeholder)
        img = Image.new("RGB", (px_width, px_height), color=(18, 18, 24))

        # Extract speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            with contextlib.suppress(Exception):
                notes_text = slide.notes_slide.notes_text_frame.text

        # Extract slide title and body text for the thumbnail
        title_text = ""
        body_texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph and ph.idx == 0:
                title_text = text
            else:
                body_texts.append(text[:200])

        # Encode to PNG base64
        buf = io.BytesIO()
        img.save(buf, format="PNG", optimize=True)
        png_b64 = base64.b64encode(buf.getvalue()).decode("ascii")

        slides_data.append({
            "index": i,
            "title": title_text,
            "body_preview": "\n".join(body_texts[:5]),
            "notes": notes_text[:2000],
            "png_base64": png_b64,
            "width": px_width,
            "height": px_height,
        })

    return JSONResponse(content={"slides": slides_data, "total": len(slides_data)})


# ---------------------------------------------------------------------------
# WebSocket — agent streaming
# ---------------------------------------------------------------------------

@app.websocket("/ws/{session_id}")
async def ws_agent(websocket: WebSocket, session_id: str) -> None:
    """Bidirectional streaming channel for a single agent session."""
    from server_adapter import (
        push_user_response,
        set_active_ws,
        set_cancel_flag,
        ws_handle_event,
        ws_reset,
        _send,
    )
    from agents import DEFAULT_TIMEOUT
    from router import route_to_agent

    await websocket.accept()
    set_active_ws(websocket)
    ws_reset()

    # Look up or create session
    session = _session_map.get(session_id)
    if session is None:
        await websocket.send_text(json.dumps(
            {"type": "error", "message": f"Session {session_id!r} not found"}
        ))
        await websocket.close()
        set_active_ws(None)
        return

    # Register WS event handler alongside existing terminal handler
    session.on(ws_handle_event)

    current_turn_task: asyncio.Task[Any] | None = None

    try:
        while True:
            raw = await websocket.receive_text()
            try:
                msg = json.loads(raw)
            except json.JSONDecodeError:
                await websocket.send_text(json.dumps(
                    {"type": "error", "message": "Invalid JSON"}
                ))
                continue

            msg_type = str(msg.get("type", ""))

            if msg_type == "message":
                content = str(msg.get("content", "")).strip()
                if not content:
                    continue

                ws_reset()

                # Determine agent via router
                try:
                    agent_name = await route_to_agent(session, content)
                except Exception:
                    agent_name = None

                if agent_name:
                    if _event_store:
                        _event_store.update_session_agent(session_id, agent_name)

                # Strip @mention prefix
                clean = content
                if content.startswith("@") and " " in content:
                    clean = content.split(" ", 1)[1]

                turn_id = None
                if _collector:
                    turn_id = _collector.on_turn_start(
                        session_id,
                        agent=agent_name or "copilot",
                        model="",
                        user_prompt=clean,
                    )

                _send({"type": "turn_started", "agent": agent_name})

                before_time = time.time()
                turn_status = "success"
                try:
                    reply = await session.send_and_wait(
                        {"prompt": clean}, timeout=DEFAULT_TIMEOUT
                    )
                except asyncio.CancelledError:
                    turn_status = "cancelled"
                    _send({"type": "cancelled"})
                except TimeoutError:
                    turn_status = "timeout"
                    _send({"type": "error", "message": "timeout"})
                except Exception as exc:
                    turn_status = "error"
                    _send({"type": "error", "message": str(exc)})

                # Detect new output files
                new_files = _find_new_outputs(before_time)
                if new_files:
                    _send({
                        "type": "new_files",
                        "files": [str(f) for f in new_files],
                    })

                _send({"type": "done", "status": turn_status})

                if _collector and turn_id:
                    _collector.on_turn_end(
                        turn_id,
                        assistant_response="",
                        status=turn_status,
                    )

            elif msg_type == "cancel":
                set_cancel_flag(True)
                if current_turn_task and not current_turn_task.done():
                    current_turn_task.cancel()

            elif msg_type == "user_response":
                content = str(msg.get("content", ""))
                push_user_response(content)

            else:
                await websocket.send_text(json.dumps(
                    {"type": "error", "message": f"Unknown message type: {msg_type}"}
                ))

    except WebSocketDisconnect:
        pass
    except Exception as exc:
        log.error("WebSocket error: %s", exc)
    finally:
        with contextlib.suppress(Exception):
            session._event_handlers = [
                h for h in session._event_handlers if h is not ws_handle_event
            ]
        set_active_ws(None)


# ---------------------------------------------------------------------------
# Output file detection (mirrors app.py logic)
# ---------------------------------------------------------------------------

_INTERESTING_SUFFIXES = {".pptx", ".md", ".py", ".bicep", ".json", ".yaml", ".sh"}


def _find_new_outputs(since: float) -> list[Path]:
    found: list[Path] = []
    grace = 3.0
    for p in _outputs_dir.rglob("*"):
        if not p.is_file():
            continue
        if any(part in _SKIP_DIRS for part in p.parts):
            continue
        if p.suffix.lower() not in _INTERESTING_SUFFIXES:
            continue
        if "-plan.md" in p.name:
            continue
        try:
            if p.stat().st_mtime >= since - grace:
                found.append(p)
        except OSError:
            pass
    return found


# ---------------------------------------------------------------------------
# Startup injection (called by server_main in app.py)
# ---------------------------------------------------------------------------

def configure(
    *,
    event_store: Any,
    copilot_client: Any,
    collector: Any,
    app_dir: Path,
    outputs_dir: Path,
) -> None:
    """Inject dependencies into the module-level singletons."""
    global _event_store, _copilot_client, _collector, _app_dir, _outputs_dir
    _event_store = event_store
    _copilot_client = copilot_client
    _collector = collector
    _app_dir = app_dir
    _outputs_dir = outputs_dir
